from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Use temporary short-named copies to avoid any character encoding issues
path_v2 = "temp_v2.pptx"
path_gen = "temp_gen.pptx"

prs_master = Presentation(path_v2)
prs_gen = Presentation(path_gen)

# 1. Add Separator Slide
# Using layout index 0 (Title) from the user's master v2
slide_layout = prs_master.slide_layouts[0]
sep_slide = prs_master.slides.add_slide(slide_layout)

# Visual polish: Match dark theme
background = sep_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(13, 21, 37)

if sep_slide.shapes.title:
    sep_slide.shapes.title.text = "TECHNICAL ARCHITECTURE & CORE MODULES"
    sep_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(103, 232, 249) # Cyan

# 2. Append Content Slides from generated presentation
# Re-creating them in the master to ensure layout consistency
for i, slide in enumerate(prs_gen.slides):
    if i < 2: continue # Skip Title and Problem slides from the generated deck
    
    title_text = slide.shapes.title.text if slide.shapes.title else "Module Analysis"
    bullet_layout = prs_master.slide_layouts[1] # Use Master's bullet layout
    new_slide = prs_master.slides.add_slide(bullet_layout)
    
    # Theme continuity
    background = new_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    # Title styling
    if new_slide.shapes.title:
        new_slide.shapes.title.text = title_text
        new_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55) # Gold
        
    # Find the content block in the source slide
    body_content = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body_content = shape.text_frame.text
            break
            
    # Inject into the new slide's main placeholder
    if body_content and len(new_slide.placeholders) > 1:
        tf = new_slide.placeholders[1].text_frame
        tf.clear()
        tf.text = body_content
        for p in tf.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(20)

merged_name = "ChemoFilter_COMPLETE_Presentation.pptx"
prs_master.save(merged_name)
print(f"Successfully created: {merged_name}")
