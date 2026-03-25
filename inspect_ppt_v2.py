from pptx import Presentation

path_v2 = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Presentation_v2.pptx"
prs = Presentation(path_v2)

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text[:100].replace('\n', ' '))
    print(" | ".join(texts))
