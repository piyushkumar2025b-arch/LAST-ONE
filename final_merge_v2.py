from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Fix: User's v2 master only has 1 slide layout! 
path_v2 = "temp_v2.pptx"
path_gen = "temp_gen.pptx"

prs_master = Presentation(path_v2)
prs_gen = Presentation(path_gen)

# 1. Add Separator Slide
# Fallback to the only layout available (index 0)
layout = prs_master.slide_layouts[0]
sep_slide = prs_master.slides.add_slide(layout)

# Theme
background = sep_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(13, 21, 37)

# Locate any shape with text on this single-layout slide to act as title
for shape in sep_slide.shapes:
    if hasattr(shape, "text_frame"):
        shape.text_frame.text = "TECHNICAL ARCHITECTURE & CORE MODULES"
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(103, 232, 249)
        break

# 2. Replicate generated slides
for i, slide in enumerate(prs_gen.slides):
    if i < 2: continue
    
    # Identify content from source
    title_text = slide.shapes.title.text if slide.shapes.title else "Module Insight"
    body_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body_text = shape.text_frame.text
            break
            
    # Create new slide in master
    new_slide = prs_master.slides.add_slide(layout)
    
    # Theme
    background = new_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    # Manually place Title and Body since we can't trust the master's single layout
    title_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf_t = title_box.text_frame
    tf_t.text = title_text
    tf_t.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55)
    tf_t.paragraphs[0].font.size = Pt(36)
    tf_t.paragraphs[0].font.bold = True

    body_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(5))
    tf_b = body_box.text_frame
    tf_b.text = body_text
    tf_b.word_wrap = True
    for p in tf_b.paragraphs:
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(20)

merged_name = "ChemoFilter_COMPLETE_Presentation.pptx"
prs_master.save(merged_name)
print(f"Successfully integrated content into: {merged_name}")
