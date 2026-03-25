from pptx import Presentation

path_v2 = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Presentation_v2.pptx"

try:
    prs_v2 = Presentation(path_v2)
    print(f"Loaded {path_v2} successfully.")
    print(f"Number of slides: {len(prs_v2.slides)}")
    
    for i, slide in enumerate(prs_v2.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        print(f"\n--- Slide {i+1}: {title} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(f"  Shape: {shape.text}")

except Exception as e:
    print(f"Error: {e}")
