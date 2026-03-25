import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Explicit paths gathered from expandproperty
path_v2 = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Presentation_v2.pptx"
path_gen = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Project_Presentation.pptx"

print(f"Checking {path_v2}: {os.path.exists(path_v2)}")
print(f"Checking {path_gen}: {os.path.exists(path_gen)}")

if not os.path.exists(path_v2) or not os.path.exists(path_gen):
    print("One or more files missing!")
    sys.exit(1)

prs_master = Presentation(path_v2)
prs_gen = Presentation(path_gen)

# Add separator slide
slide_layout = prs_master.slide_layouts[0]
sep_slide = prs_master.slides.add_slide(slide_layout)
background = sep_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(13, 21, 37)

if sep_slide.shapes.title:
    sep_slide.shapes.title.text = "TECHNICAL ARCHITECTURE & CORE MODULES"
    sep_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(103, 232, 249)

# Now manually add the content slides from prs_gen
for i, slide in enumerate(prs_gen.slides):
    if i < 2: continue  # Skip title and problem statement if already in v2
    
    title_text = slide.shapes.title.text if slide.shapes.title else "Module Details"
    bullet_slide_layout = prs_master.slide_layouts[1]
    new_slide = prs_master.slides.add_slide(bullet_slide_layout)
    
    background = new_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    if new_slide.shapes.title:
        new_slide.shapes.title.text = title_text
        new_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55)
    
    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body_shape = shape
            break
            
    if body_shape and len(new_slide.placeholders) > 1:
        tf = new_slide.placeholders[1].text_frame
        tf.clear() # Clear placeholder default text
        tf.text = body_shape.text_frame.text
        for p in tf.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(22)

final_path = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_FINAL_Merged_Deck.pptx"
prs_master.save(final_path)
print(f"Succesfully merged into: {final_path}")
