import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Fix filenames based on inspection and target
path_v2 = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Presentation_v2.pptx"
path_gen = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Project_Presentation.pptx"

if not os.path.exists(path_v2):
    print(f"Error: Could not find {path_v2}")
    sys.exit(1)
if not os.path.exists(path_gen):
    print(f"Error: Could not find {path_gen}")
    sys.exit(1)

prs_master = Presentation(path_v2)
prs_gen = Presentation(path_gen)

# Add separator slide
slide_layout = prs_master.slide_layouts[0]
sep_slide = prs_master.slides.add_slide(slide_layout)
background = sep_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(13, 21, 37)

if sep_slide.shapes.title:
    sep_slide.shapes.title.text = "TECHNICAL ARCHITECTURE & CORE MODULES"
    sep_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(103, 232, 249)

# Now manually add the content slides from prs_gen
# Slide 2 in gen is "Problem Statement", Slide 3 is "Project Overview"
# Let's start from index 2 (Slide 3)
for i, slide in enumerate(prs_gen.slides):
    if i < 2: continue  # Skip title and problem
    
    title_text = slide.shapes.title.text if slide.shapes.title else "Module Details"
    bullet_slide_layout = prs_master.slide_layouts[1]
    new_slide = prs_master.slides.add_slide(bullet_slide_layout)
    
    # Custom Background
    background = new_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    if new_slide.shapes.title:
        new_slide.shapes.title.text = title_text
        new_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55)
    
    # Copy bullets
    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body_shape = shape
            break
            
    if body_shape and len(new_slide.placeholders) > 1:
        tf = new_slide.placeholders[1].text_frame
        tf.text = body_shape.text_frame.text
        for p in tf.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(22)

final_path = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_FINAL_Merged_Deck.pptx"
prs_master.save(final_path)
print(f"Succesfully merged into: {final_path}")
