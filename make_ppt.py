from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Use standard 16:9 aspect ratio if possible, otherwise stick to default (4:3)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Helper functions ---
def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Styling
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(103, 232, 249) # Cyan
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    return slide

def add_bullet_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55) # Gold
    title.text_frame.paragraphs[0].font.name = "Arial"
    
    tf = body.text_frame
    tf.text = bullet_points[0]
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].font.size = Pt(24)
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(24)
        p.level = 0
        
    return slide

# --- Build the Presentation ---

# Slide 1: Title
add_title_slide(
    "ChemoFilter: Advanced Computational Drug Discovery Platform",
    "Streamlining in-silico screening, target prioritization, and molecular validation."
)

# Slide 2: Problem Statement
add_bullet_slide(
    "The Structural Bottleneck in Early Drug Discovery",
    [
        "High Failure Rates: 90% of clinical candidates fail due to unseen toxicity or poor ADME profiles.",
        "Data Fragmentation: Computational chemists traditionally jump between dozens of isolated tools (PAINS filters, logP calculators, 3D modeling).",
        "Slow Assessment: Manual inspection of SMILES strings and topological metrics takes months of human-hours.",
        "The Need: A unified, intelligent, ultra-fast platform to filter millions of candidates into high-quality hits."
    ]
)

# Slide 3: What is ChemoFilter?
add_bullet_slide(
    "ChemoFilter Platform Overview",
    [
        "An AI-powered, multi-tab computational chemistry web application built perfectly on Streamlit and RDKit.",
        "Instantly analyzes molecular representations (SMILES) to predict drug-likeness, safety, and synthetic accessibility.",
        "Architecture: Operates entirely locally with robust Python backends, avoiding cloud exposure of proprietary IP.",
        "Scalability: Equipped with Mega-Iterative 'Zenith' and 'Celestial' evaluation engines checking up to 50+ physiochemical parameters in milliseconds."
    ]
)

# Slide 4: Core Capabilities - The Analysis Engine
add_bullet_slide(
    "Pillar 1: Deep ADMET & Physicochemical Analysis",
    [
        "Lipinski's Rule of 5 Validation (MW, LogP, HBD, HBA) for oral bioavailability.",
        "Advanced topological scoring: QED (Quantitative Estimate of Drug-likeness), TPSA, and Fsp3 orbital hybridization.",
        "3D Conformational Stress calculations via the MMFF force field.",
        "Automated Toxicity Alerts: Checks structures against SMARTS libraries for PAINS (Pan Assay Interference Compounds), Brenk alerts, Ames mutagenicity, and hERG inhibition."
    ]
)

# Slide 5: Core Capabilities - Visualization & AI
add_bullet_slide(
    "Pillar 2: Data Visualization & AI Explanations",
    [
        "Rich 2D & 3D Visuals: Interactive 3D molecular conformity viewers, Plotly radar charts, and scaffold clustering graphs.",
        "LLM Integration: Anthropic Claude integration acts as a virtual medicinal chemist.",
        "Generates analog suggestions, identifies synthetic vulnerabilities, and predicts biological mechanisms instantly from the evaluated metrics.",
        "Side-by-side Gold Standard tracking allows direct comparison against market reference compounds like Olanzapine."
    ]
)

# Slide 6: The Scoring System
add_bullet_slide(
    "ChemoScore: The Ultimate Prioritization Metric",
    [
        "Weighted Evaluation: Calculates an overall 'ChemoScore' and Grade (A, B, C, F) for any input.",
        "Adjustable Sliders: Lead reviewers can dynamically distribute weight among Structure, PhysChem, Safety, and Synthesis modules.",
        "Ensures researchers aren't just optimizing for binding affinity, but prioritizing molecules that can actually survive the human liver and bloodstream.",
        "Actionable outputs: Allows instant sorting of thousands of hits down to the Top 5 most viable structures."
    ]
)

# Slide 7: Technical Stack & Innovation
add_bullet_slide(
    "Technical Architecture & Robustness",
    [
        "Frontend: Streamlit configured with custom-injected HTML/CSS for a dark, futuristic laboratory interface.",
        "Backend: Python, RDKit, Pandas, and NetworkX. Highly localized to ensure zero molecular data leakage.",
        "Caching Optimization: Utilizes advanced session caching and @st.cache_resource decorators to stabilize execution and prevent memory serialization crashes.",
        "Result: Real-time calculation speeds previously only achievable by massive compute clusters."
    ]
)

# Slide 8: Future Scope & Impact
add_bullet_slide(
    "Platform Impact & Next Steps",
    [
        "Impact: Decreases hit-to-lead verification time from weeks to seconds.",
        "Significantly de-risks preclinical assets by catching structural red flags (Cytochrome P450 liabilities, false positives) early.",
        "Next Steps: Integration with broader high-throughput screening data lakes.",
        "Expansion of generative AI capabilities to auto-design synthesis pathways."
    ]
)

# Slide 9: Q&A
add_title_slide(
    "Thank You",
    "Ready for Live Demonstration & Questions"
)

prs.save("ChemoFilter_Project_Presentation.pptx")
print("Presentation generated successfully at ChemoFilter_Project_Presentation.pptx")
