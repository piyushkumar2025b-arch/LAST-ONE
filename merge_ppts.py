from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load both presentations
path_v2 = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Presentation_v2.pptx"
path_gen = r"C:\Users\Piyush Kumar\OneDrive\Attachments\zip MDP\LAST ONE\ChemoFilter_Project_Presentation.pptx"

# We will use v2 as the base to keep the user's styling/master
prs_master = Presentation(path_v2)
prs_gen = Presentation(path_gen)

def copy_slide_content(source_slide, target_slide):
    # This is a basic copy function since python-pptx doesn't have native slide merging
    for shape in source_slide.shapes:
        if shape.has_text_frame:
            # We add a text box at the same position
            new_shape = target_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text_frame.text = shape.text_frame.text
            # Basic style copy for the first paragraph
            if len(shape.text_frame.paragraphs) > 0 and len(new_shape.text_frame.paragraphs) > 0:
                 new_shape.text_frame.paragraphs[0].font.size = shape.text_frame.paragraphs[0].font.size
                 new_shape.text_frame.paragraphs[0].font.color.rgb = shape.text_frame.paragraphs[0].font.color.rgb

# Since v2 seems to be the UX/UI heavy one, and my generated one is content heavy,
# I will append the generated content slides after the relevant sections in v2.
# However, without exact layouts, the safest is to append the new slides at the end 
# as a "Detailed Technical Appendix" or "Advanced Modules" section.

# Add a separator slide in prs_master
slide_layout = prs_master.slide_layouts[0] # Title slide layout
sep_slide = prs_master.slides.add_slide(slide_layout)
# Set background
background = sep_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(13, 21, 37)

if sep_slide.shapes.title:
    sep_slide.shapes.title.text = "TECHNICAL ARCHITECTURE & CORE MODULES"
    sep_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(103, 232, 249)

# Now manually add the content slides from prs_gen (skip title/problem if redundant)
# For better reliability, we will recreate the slides to match the master's aspect ratio
for i, slide in enumerate(prs_gen.slides):
    if i == 0 or i == 1: continue # Skip title and problem statement as v2 likely has them
    
    title_text = slide.shapes.title.text if slide.shapes.title else "Module Details"
    bullet_slide_layout = prs_master.slide_layouts[1] # Bullet layout
    new_slide = prs_master.slides.add_slide(bullet_slide_layout)
    
    # Custom Background
    background = new_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 21, 37)
    
    if new_slide.shapes.title:
        new_slide.shapes.title.text = title_text
        new_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(212, 175, 55)
    
    # Copy bullets
    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body_shape = shape
            break
            
    if body_shape and len(new_slide.placeholders) > 1:
        tf = new_slide.placeholders[1].text_frame
        tf.text = body_shape.text_frame.text
        for p in tf.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(22)

prs_master.save("ChemoFilter_Final_Merged_Presentation.pptx")
print("Merged presentation created: ChemoFilter_Final_Merged_Presentation.pptx")
