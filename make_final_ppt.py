import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Colors
    bg_color = RGBColor(10, 14, 26)  # Midnight Navy #0a0e1a
    accent_color = RGBColor(232, 160, 32)  # Amber Gold #E8A020
    text_color = RGBColor(255, 255, 255)  # White
    gray_color = RGBColor(200, 200, 200)  # Light Gray

    def apply_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_styled_title(slide, title_text):
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].font.size = Pt(36)
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.color.rgb = accent_color
        title_text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add underline
        left = Inches(0.5)
        top = Inches(1.1)
        width = Inches(9)
        height = Inches(0.05)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_color
        shape.line.visible = False

    def add_bullet_points(slide, points, font_size=Pt(20)):
        if not slide.placeholders[1]:
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
        else:
            content_shape = slide.placeholders[1]
            
        tf = content_shape.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = font_size
            p.font.color.rgb = text_color
            p.space_after = Pt(10)
            p.level = 0

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_bg(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ChemoFilter: AI-Powered Drug Discovery Platform"
    subtitle.text = "Crystalline Noir Edition — Omnipotent v1,000,000\nVIT Chennai MDP 2026 | March 2026"
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    subtitle.text_frame.paragraphs[0].font.color.rgb = gray_color

    # 2. The Drug Discovery Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "The Drug Discovery Problem")
    add_bullet_points(slide, [
        "90% Failure Rate: Most drug candidates fail in clinical trials due to toxicity or poor ADMET.",
        "Staggering Costs: Average cost to bring a new drug to market is $2.6 Billion.",
        "Time Intensive: Development cycle takes 12-15 years on average.",
        "The ADMET Bottleneck: Failure to identify issues early leads to massive late-stage losses.",
        "Need for In Silico Screening: Computational filters can prioritize safer, more effective leads."
    ])

    # 3. What is ChemoFilter?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "What is ChemoFilter?")
    add_bullet_points(slide, [
        "Comprehensive Platform: A multi-parameter ADMET screening and drug-likeness engine.",
        "Target Audience: Medicinal chemists, researchers, and computational students.",
        "Mission: Accelerate lead optimization through intelligent filtering and AI-driven insights.",
        "Crystalline Noir Aesthetic: High-performance UI designed for professional scientific environments.",
        "Context: Developed for VIT Chennai Management Development Program 2026."
    ])

    # 4. UI Layout Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "User Interface & Design")
    add_bullet_points(slide, [
        "Midnight Navy & Amber Gold: Aesthetic inspired by high-end scientific instrumentation.",
        "Interactive Sidebar: Central control for SMILES input, file uploads, and engine selection.",
        "Dynamic Stats Strip: Instant visual feedback on dataset quality and property counts.",
        "Live Leaderboard: Ranked compound dossiers with real-time sorting and filtering.",
        "Integrated FAB System: Floating panels for deep technical data and reference vaults."
    ])

    # 5. Key Features Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Key Features Overview")
    add_bullet_points(slide, [
        "21+ Advanced Descriptors: Exhaustive Physicochemical and ADMET profiling.",
        "Multi-Tox Screening: PAINS, Brenk, hERG, Ames, and covalent warhead detection.",
        "Integrated AI: Scientific explanation narratives powered by Gemini / Claude.",
        "Visualization Suite: Radar charts, BOILED-Egg plots, and similarity heatmaps.",
        "Batch Intelligence: Parallel processing of libraries up to 200 compounds."
    ])

    # 6. ADMET Screening Explained
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "ADMET: The Core Science")
    add_bullet_points(slide, [
        "Absorption: HIA, Caco-2, and ESOL solubility modelling.",
        "Distribution: BBB penetration, Plasma Protein Binding (PPB), and LogD predictions.",
        "Metabolism: 5-Isoform CYP450 panel (3A4, 2D6, 2C9, 2C19, 1A2).",
        "Excretion: Renal clearance estimation and molecular half-life flagging.",
        "Toxicity: Multi-layer risk modeling for cardiac, mutagenic, and liver health."
    ])

    # 7. Lipinski Rule of Five
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Lipinski Rule of Five & Filters")
    add_bullet_points(slide, [
        "Rule of Five (Ro5): Foundation of oral drug-likeness assessment.",
        "Key Thresholds: MW < 500, LogP < 5, HBD < 5, HBA < 10.",
        "Extended Filters: Veber (TPSA/RotB), Ghose, Egan, and Muegge integration.",
        "Flexible Compliance: Automated detection of violations and penalty scoring.",
        "Optimization Focus: Identifying 'lead-like' vs 'drug-like' space."
    ])

    # 8. QED and SA Score
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "QED & Synthetic Accessibility")
    add_bullet_points(slide, [
        "QED (Quantitative Estimate of Drug-likeness): Integrated desirability function (0.0 to 1.0).",
        "SA Score: Ertl's fragment-based synthetic ease prediction (1.0 = easy, 10.0 = hard).",
        "Multi-Param Optimization: Balancing chemistry complexity with drug potential.",
        "Desirability Mapping: Prioritizing compounds that hit the 'sweet spot' for all properties.",
        "Natural Product Likeness: NP-Score filtering for discovery inspiration."
    ])

    # 9. Safety & Structural Alerts
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Safety & Structural Alerts")
    add_bullet_points(slide, [
        "PAINS Filter: 480+ patterns for Pan-Assay Interference Compounds.",
        "Brenk Filter: Detection of reactive or undesirable functional groups.",
        "hERG Risk: Cardiac liability prediction to prevent QT prolongation.",
        "Ames Mutagenicity: Virtual Salmonella test for DNA damage risk.",
        "Covalent Warheads: Identifying potential irreversible binders early."
    ])

    # 10. The Multi-Engine Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "The Omnipotent Engine Pipeline")
    add_bullet_points(slide, [
        "Layered Tier System: 10 computational engines from Vanguard to Aether.",
        "Modular Routing: Orchestrator dynamically shifts SMILES through tier complexity.",
        "Vanguard Core (v2.0): High-speed initial screening and descriptor compute.",
        "Zenith Engines: Deep organ toxicity, metabolic site prediction, and costs.",
        "Aether-Primality (v10000): Ultimate analysis tier for God-mode final reports."
    ])

    # 11. Scoring System: ChemoScore™
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Scoring System: ChemoScore™")
    add_bullet_points(slide, [
        "Integrity (20%): Structural validity, organic checks, and valency.",
        "PhysChem (25%): Optimal MW, LogP, and TPSA balance.",
        "Potency/Discovery (25%): QED and Ligand Efficiency.",
        "Safety (20%): Toxicophore alerts and hERG/Ames risk penalization.",
        "Synthesis (10%): SA Score and complexity thresholds.",
        "Grades: A (80+), B (60+), C (40+), F (<40)."
    ])

    # 12. AI Integration & Explainer
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "AI Integration: Gemini Powered")
    add_bullet_points(slide, [
        "Narrative Analysis: AI translates raw ADMET numbers into clear reports.",
        "Three Modes: Overview, Toxicity Deep Dive, and Optimization Suggestions.",
        "Strategic Guidance: AI suggests medicinal chemistry modifications to fix flags.",
        "Unified Interface: Toggle between technical stats and plain-language summaries.",
        "Offline Stability: Framework continues to work 100% even if AI key is missing."
    ])

    # 13. Batch Processing & Synthetic Data
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Batch Intelligence")
    add_bullet_points(slide, [
        "Library Throughput: Parallel processing of large SMILES collections.",
        "Synthetic Filler: Auto-generates dummy compounds for UI stress testing.",
        "Data Hygiene: Recursive deduplication and SMILES normalization.",
        "Session Persistence: All batch results cached for instant tab-switching.",
        "Dataset Stats: Real-time calculation of histograms and property distributions."
    ])

    # 14. Visualization Suite
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Visualization & Analytics")
    add_bullet_points(slide, [
        "Radar QED Maps: Octagonal profile of compound drug-likeness.",
        "BOILED-Egg: Human Intestinal Absorption vs BBB penetration scatter plot.",
        "Pulse Histograms: Population distribution of MW, LogP, and SA Scores.",
        "Parallel Coordinates: Multi-parametric filtering across all descriptors.",
        "Interactive Heatmaps: Color-coded grids for grade comparisons."
    ])

    # 15. Input & Data Methods
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Input & Connectivity")
    add_bullet_points(slide, [
        "Direct SMILES: Comma-separated strings with live validation.",
        "CSV/Excel Uploads: Bulk import of thousands of candidates.",
        "PubChem Lookup: Name-to-SMILES resolution via live API integration.",
        "Export Control: Download reports in CSV, PDF-ready HTML, or JSON.",
        "Reference Benchmarks: Comparing leads directly against FDA approved sets."
    ])

    # 16. Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Modern Technology Stack")
    add_bullet_points(slide, [
        "Core Language: Python 3.12 (High-performance scripting).",
        "Cheminformatics: RDKit (Descriptors, SMARTS, 3D conformation).",
        "Framework: Streamlit (Reactive UI and session management).",
        "Data Engine: Pandas & NumPy (Tensor ops and vectorized matching).",
        "AI: Google Gemini Pro / Anthropic Claude Integration."
    ])

    # 17. Scientific References
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Scientific Foundation")
    add_bullet_points(slide, [
        "[1] Daina A, et al. BOILED-Egg (HIA/BBB), ChemMedChem 2016.",
        "[2] Lipinski C, et al. Rule of Five, ADDR 2001.",
        "[3] Delaney J, et al. ESOL Solubility Model, JCICS 2004.",
        "[4] Bickerton G, et al. QED Desirability, Nat Chem 2012.",
        "[5] Baell J, et al. PAINS Filters, J Med Chem 2010.",
        "[6] Ertl P, et al. Synthetic Accessibility (SA Score), J Cheminf 2009."
    ])

    # 18. Deployment & Infrastructure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Cloud Deployment")
    add_bullet_points(slide, [
        "Streamlit Cloud: Auto-deployment synced with GitHub repository.",
        "Secrets Management: Secure storage of AI and API credentials.",
        "Packages.txt: Automated installation of RDKit OS-level dependencies.",
        "Modular Cache: st.cache_data for 24h results persistence.",
        "Resource Efficient: Runs on standard cloud instances with low footprint."
    ])

    # 19. Results & Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Project Results")
    add_bullet_points(slide, [
        "Reliability: Zero-crash stability across all tiered engines.",
        "Accuracy: Validated descriptor calculation against RDKit standards.",
        "User Experience: Seamless transition from landing to full reports.",
        "Integration: Successful merging of physics, chemistry, and AI layers.",
        "Outcome: A professional-grade sandbox for initial drug screening."
    ])

    # 20. Future Directions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "The Road Ahead")
    add_bullet_points(slide, [
        "3D Docking: In-browser protein-ligand interaction visualization.",
        "Advanced QSAR: Machine learning models for specific target affinity.",
        "Chemical Search: Direct similarity searching against 100M+ PubChem items.",
        "Mobile App: Responsive PWA for field researcher access.",
        "API Service: REST endpoints for integration into pharma workflows."
    ])

    # 21. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Conclusion")
    add_bullet_points(slide, [
        "ChemoFilter transforms raw SMILES into actionable drug intelligence.",
        "Combines classical chemistry rules with cutting-edge Generative AI.",
        "Built for speed, accuracy, and beauty.",
        "Version 1,000,000 is deployment-ready.",
        "Empowering the next generation of medicinal chemists."
    ])

    # 22. Thank You / Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_bg(slide)
    add_styled_title(slide, "Thank You")
    add_bullet_points(slide, [
        "VIT Chennai MDP 2026 Presentation",
        "Author: Piyush Kumar (25BCE1392)",
        "GitHub: [Link to Project Repository]",
        "Acknowledgements: RDKit, Streamlit, and Google AI Teams.",
        "Questions? — Open Floor"
    ])

    prs.save('ChemoFilter_Final_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()
